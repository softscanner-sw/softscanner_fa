"""
Generate scientific presentation for SoftScanner FA project.
Sources: main.pdf, approach.md, subjects.md, autoe2e-benchmark-evaluation-report.md
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
MEDIUM_BLUE = RGBColor(0x2C, 0x5F, 0x8A)
LIGHT_BLUE = RGBColor(0x3A, 0x7C, 0xBD)
ACCENT_ORANGE = RGBColor(0xE8, 0x6C, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MED_GRAY = RGBColor(0x66, 0x66, 0x66)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
RED = RGBColor(0xC6, 0x28, 0x28)
AMBER = RGBColor(0xF5, 0x7F, 0x17)

def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=DARK_GRAY, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = spacing
        p.level = 0
    return txBox

def add_table(slide, left, top, width, height, rows, cols, data, col_widths=None, header_color=DARK_BLUE, font_size=12):
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c]) if r < len(data) and c < len(data[r]) else ""
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.name = 'Calibri'
                if r == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                else:
                    paragraph.font.color.rgb = DARK_GRAY
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
    return table

def section_header(slide, text):
    add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), DARK_BLUE)
    add_text_box(slide, Inches(0.6), Inches(0.2), Inches(12), Inches(0.8), text, 32, True, WHITE)

# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, prs.slide_height, DARK_BLUE)
add_shape_bg(slide, Inches(0), Inches(5.5), prs.slide_width, Inches(0.06), ACCENT_ORANGE)

add_text_box(slide, Inches(1), Inches(1.5), Inches(11.3), Inches(1.5),
    "Automating End-to-End Functional Testing\nof Web Applications based on\nAugmented Workflows and LLMs",
    36, True, WHITE, PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(4.0), Inches(11.3), Inches(0.5),
    "EASE 2026 — 30th International Conference on Evaluation and Assessment in Software Engineering",
    18, False, RGBColor(0xAA, 0xCC, 0xEE), PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(4.7), Inches(11.3), Inches(0.5),
    "Glasgow, United Kingdom  |  9–12 June 2026",
    16, False, RGBColor(0x88, 0xAA, 0xCC), PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(6.0), Inches(11.3), Inches(0.5),
    "Current implementation status through Phase B2 — validated on 6 Angular applications",
    14, False, RGBColor(0x88, 0xAA, 0xCC), PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2: Problem Context & Motivation
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_header(slide, "Problem Context and Motivation")

add_text_box(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(0.5),
    "E2E functional testing validates that frontend behavior works correctly from the user's perspective.",
    18, True, DARK_BLUE)

items = [
    "User-visible behavior is primarily defined at the frontend — defects at this level are invisible to unit/backend testing",
    "E2E testing proceeds by defining workflows: structured sequences of user actions (clicks, inputs, navigations)",
    "Executable tests are concrete instantiations of such workflows — coverage depends on how workflows are defined",
    "Existing approaches compute coverage over interaction spaces derived from models, specifications, or runtime artifacts",
    "Coverage remains bounded by abstraction choices or exploration policies, not by the implementation itself",
]
add_bullet_list(slide, Inches(0.8), Inches(2.0), Inches(5.8), Inches(4.5), items, 15)

# Right column: the gap
add_shape_bg(slide, Inches(7.0), Inches(1.8), Inches(5.8), Inches(4.8), RGBColor(0xFD, 0xF0, 0xE0))
add_text_box(slide, Inches(7.2), Inches(1.9), Inches(5.4), Inches(0.5),
    "The Coverage Gap", 20, True, ACCENT_ORANGE)

gap_items = [
    "Static models abstract away guards, roles, parameters, and state constraints",
    "Dynamic exploration is bounded by exploration budget and observed behavior",
    "Neither constructs a finite, implementation-grounded workflow reference space",
    "Result: coverage reflects how well an approach explores its chosen artifact, not how completely it exercises the implementation",
]
add_bullet_list(slide, Inches(7.2), Inches(2.5), Inches(5.4), Inches(3.8), gap_items, 14, DARK_GRAY)

# ============================================================
# SLIDE 3: Research Objective & Questions
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_header(slide, "Research Objective and Questions")

add_text_box(slide, Inches(0.6), Inches(1.4), Inches(12), Inches(0.8),
    "Objective: Maximize E2E test coverage over a finite, constraint-aware, implementation-grounded workflow space.",
    20, True, DARK_BLUE)

# RQ boxes
rq_data = [
    ("RQ1 — Representation", "How can frontend interaction behavior be represented such that the resulting workflow space is complete w.r.t. the implementation, constraint-aware, finite, and executable?",
     "Representational inadequacy: classical workflows capture reachability but not executability"),
    ("RQ2 — Construction", "How can such a representation be systematically and deterministically derived from frontend source code?",
     "Methodological: even with a suitable representation, constructing it from code must be systematic, deterministic, and finite"),
    ("RQ3 — Realization", "How can workflows in this representation be instantiated and executed so as to maximize coverage over that space?",
     "Operational: feasibility does not guarantee execution success — realization requires concrete runtime assignments"),
]

for i, (title, question, issue) in enumerate(rq_data):
    y = Inches(2.5) + Inches(i * 1.6)
    add_shape_bg(slide, Inches(0.6), y, Inches(12), Inches(1.4), LIGHT_GRAY if i % 2 == 0 else WHITE)
    add_text_box(slide, Inches(0.8), y + Inches(0.05), Inches(3), Inches(0.4), title, 17, True, MEDIUM_BLUE)
    add_text_box(slide, Inches(0.8), y + Inches(0.45), Inches(11.5), Inches(0.4), question, 14, False, DARK_GRAY)
    add_text_box(slide, Inches(0.8), y + Inches(0.9), Inches(11.5), Inches(0.4), "Issue: " + issue, 13, False, MED_GRAY)


# ============================================================
# SLIDE 4: Related Work — Why Existing Approaches Are Insufficient
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_header(slide, "Why Existing Automated E2E Generation Is Insufficient")

# Left: Static
add_shape_bg(slide, Inches(0.4), Inches(1.3), Inches(6.0), Inches(5.6), LIGHT_GRAY)
add_text_box(slide, Inches(0.6), Inches(1.4), Inches(5.6), Inches(0.5),
    "Static Construction", 20, True, MEDIUM_BLUE)
items_static = [
    "Model-Based Testing (MBT): finite state machines, state/transition coverage",
    "Use case / requirements-driven: UML models, textual specs",
    "LLM-based synthesis from user stories, Gherkin, form descriptions",
    "Limitation: interaction space defined by the specification, not the implementation",
    "Generated scripts may target unreachable states, miss guards/constraints",
    "Coverage bounded by artifact completeness and alignment with implementation",
]
add_bullet_list(slide, Inches(0.6), Inches(2.0), Inches(5.6), Inches(4.5), items_static, 14)

# Right: Dynamic
add_shape_bg(slide, Inches(6.8), Inches(1.3), Inches(6.0), Inches(5.6), RGBColor(0xE8, 0xF0, 0xFE))
add_text_box(slide, Inches(7.0), Inches(1.4), Inches(5.6), Inches(0.5),
    "Dynamic Exploration", 20, True, MEDIUM_BLUE)
items_dynamic = [
    "Crawling: AJAX Crawling, CrawlJax, GUI ripping — discover observed DOM states",
    "Web agents: LLMs operating on live DOM, selecting actions step-by-step",
    "AutoE2E: LLM-driven feature discovery from runtime exploration",
    "Limitation: interaction space defined by exploration budget and observed behavior",
    "May miss feasible but unvisited workflows; metrics depend on exploration strategy",
    "Coverage reflects discovered subset, not the full implementation-level workflow space",
]
add_bullet_list(slide, Inches(7.0), Inches(2.0), Inches(5.6), Inches(4.5), items_dynamic, 14)

# Bottom banner
add_shape_bg(slide, Inches(0.4), Inches(6.2), Inches(12.5), Inches(0.6), ACCENT_ORANGE)
add_text_box(slide, Inches(0.6), Inches(6.25), Inches(12), Inches(0.5),
    "Gap: No prior work constructs a finite, constraint-aware workflow reference space grounded directly in the frontend implementation.",
    15, True, WHITE, PP_ALIGN.CENTER)

# ============================================================
# SLIDE 5: Core Idea
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_header(slide, "Core Idea: Augmented Workflows over a Closed Workflow Space")

add_text_box(slide, Inches(0.6), Inches(1.4), Inches(12), Inches(0.5),
    "Key Insight: Separate representation construction (Phase A) from realization (Phase B).",
    20, True, DARK_BLUE)

# Augmented workflow definition
add_shape_bg(slide, Inches(0.6), Inches(2.2), Inches(5.8), Inches(4.5), LIGHT_GRAY)
add_text_box(slide, Inches(0.8), Inches(2.3), Inches(5.4), Inches(0.5),
    "Augmented Workflows", 18, True, MEDIUM_BLUE)
items_aw = [
    "A workflow w is a path whose transitions are annotated with constraints extracted from the implementation",
    "Constraints originate from route guards, route metadata (roles), required parameters, and atomic UI predicates",
    "All constraints are aggregated into a workflow-level constraint set C(w)",
    "w is FEASIBLE if C(w) is not contradictory; CONDITIONAL if it requires runtime assignment; PRUNED if contradictory",
    "The reference space W = all non-PRUNED workflows — fixed, finite, implementation-grounded",
]
add_bullet_list(slide, Inches(0.8), Inches(2.9), Inches(5.4), Inches(3.5), items_aw, 14)

# Coverage definition
add_shape_bg(slide, Inches(6.8), Inches(2.2), Inches(5.8), Inches(4.5), RGBColor(0xE8, 0xF0, 0xFE))
add_text_box(slide, Inches(7.0), Inches(2.3), Inches(5.4), Inches(0.5),
    "Coverage over W", 18, True, MEDIUM_BLUE)
items_cov = [
    "W is constructed once (Phase A) and held fixed",
    "Phase B realizes workflows in W into executable Selenium tests via validated assignments",
    "Coverage = |E| / |W| where E = successfully executed workflows",
    "Both our approach and baselines are evaluated over the same W",
    "Coverage differences reflect realization effectiveness, not denominator differences",
    "This separates what the implementation can do from how well an approach exercises it",
]
add_bullet_list(slide, Inches(7.0), Inches(2.9), Inches(5.4), Inches(3.5), items_cov, 14)

# ============================================================
# SLIDE 6: Pipeline Overview
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_header(slide, "End-to-End Pipeline: Phase A (Static) + Phase B (Realization)")

# Phase A header
add_shape_bg(slide, Inches(0.4), Inches(1.3), Inches(6.2), Inches(0.5), MEDIUM_BLUE)
add_text_box(slide, Inches(0.5), Inches(1.33), Inches(6), Inches(0.4),
    "Phase A — Static Construction of Workflow Space W", 16, True, WHITE)

# A1 box
add_shape_bg(slide, Inches(0.4), Inches(1.9), Inches(2.8), Inches(2.4), RGBColor(0xDE, 0xEA, 0xF7))
add_text_box(slide, Inches(0.5), Inches(1.95), Inches(2.6), Inches(0.35), "A1: Multigraph Extraction", 14, True, DARK_BLUE)
a1_items = ["Angular AST \u2192 multigraph G", "6 node kinds, 18 edge kinds", "Constraint metadata on edges", "Deterministic, auditable"]
add_bullet_list(slide, Inches(0.5), Inches(2.35), Inches(2.6), Inches(1.8), a1_items, 12)

# arrow
add_text_box(slide, Inches(3.25), Inches(2.8), Inches(0.5), Inches(0.5), "\u2192", 28, True, ACCENT_ORANGE, PP_ALIGN.CENTER)

# A2 box
add_shape_bg(slide, Inches(3.6), Inches(1.9), Inches(2.8), Inches(2.4), RGBColor(0xDE, 0xEA, 0xF7))
add_text_box(slide, Inches(3.7), Inches(1.95), Inches(2.6), Inches(0.35), "A2: Workflow Enumeration", 14, True, DARK_BLUE)
a2_items = ["1 workflow per trigger edge", "Handler-scoped effect closure", "Constraint merge + classify", "W = {FEASIBLE \u222a CONDITIONAL}"]
add_bullet_list(slide, Inches(3.7), Inches(2.35), Inches(2.6), Inches(1.8), a2_items, 12)

# Phase B header
add_shape_bg(slide, Inches(0.4), Inches(4.5), Inches(12.5), Inches(0.5), ACCENT_ORANGE)
add_text_box(slide, Inches(0.5), Inches(4.53), Inches(12), Inches(0.4),
    "Phase B — Coverage-Guided Realization over W", 16, True, WHITE)

# B0-B4 boxes
b_stages = [
    ("B0", "Subject\nManifest", "Manual config:\ncredentials,\nroute params,\nform data"),
    ("B1", "Intent +\nActionPlan", "Deterministic\nderivation:\nassignments,\npreconditions"),
    ("B2", "Selenium\nCode Gen", "1 test file per\nActionPlan,\ndeterministic"),
    ("B3", "Execution\n+ Retry", "Bounded retry\n(3 levels),\nfailure taxonomy"),
    ("B4", "Coverage\nReporting", "C1-C4 metrics\nover fixed W"),
]
for i, (label, title, desc) in enumerate(b_stages):
    x = Inches(0.4 + i * 2.55)
    add_shape_bg(slide, x, Inches(5.15), Inches(2.3), Inches(2.0), RGBColor(0xFD, 0xF0, 0xE0))
    add_text_box(slide, x + Inches(0.05), Inches(5.2), Inches(2.2), Inches(0.35), f"{label}: {title}", 13, True, ACCENT_ORANGE)
    add_text_box(slide, x + Inches(0.05), Inches(5.6), Inches(2.2), Inches(1.4), desc, 11, False, DARK_GRAY)
    if i < 4:
        add_text_box(slide, x + Inches(2.3), Inches(5.9), Inches(0.3), Inches(0.4), "\u2192", 20, True, MED_GRAY, PP_ALIGN.CENTER)

# Status indicators in Phase A/B pipeline
add_shape_bg(slide, Inches(6.8), Inches(1.3), Inches(6.0), Inches(3.0), RGBColor(0xE8, 0xF5, 0xE9))
add_text_box(slide, Inches(7.0), Inches(1.4), Inches(5.6), Inches(0.4), "Current Validated Status", 16, True, GREEN)
status_items = [
    "\u2713 A1: Multigraph extraction (6 subjects)",
    "\u2713 A2: 257 workflows enumerated + classified",
    "\u2713 B0: 6 manifests validated",
    "\u2713 B1: 257 intents + 257 plans (GT-aligned)",
    "\u2713 B2: 257 Selenium tests generated (100%)",
    "\u25cb B3: Execution — not yet validated at scale",
    "\u25cb B4: Coverage reporting — pending B3",
]
add_bullet_list(slide, Inches(7.0), Inches(1.9), Inches(5.6), Inches(2.2), status_items, 13, DARK_GRAY)


# ============================================================
# SLIDE 7: A1 — UI Interaction Multigraph
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_header(slide, "A1: UI Interaction Multigraph (RQ1, RQ2)")

add_text_box(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(0.5),
    "A1 statically analyzes the Angular codebase via AST inspection to construct a unified multigraph G.",
    17, False, DARK_GRAY)

# Node kinds table
node_data = [
    ["Node Kind", "Represents", "Identity"],
    ["Module", "Angular module / standalone root", "file + class name"],
    ["Route", "URL-level route context", "canonical fullPath"],
    ["Component", "Angular component class", "file + class name"],
    ["Widget", "Interactive template element", "component + span + kind"],
    ["Service", "Injectable service class", "file + class name"],
    ["External", "External URL target", "FNV-1a hash of URL"],
]
add_table(slide, Inches(0.4), Inches(1.9), Inches(6.0), Inches(3.0), 7, 3, node_data,
    [Inches(1.2), Inches(2.8), Inches(2.0)], font_size=12)

# Edge kinds summary
add_shape_bg(slide, Inches(6.8), Inches(1.9), Inches(6.0), Inches(5.0), LIGHT_GRAY)
add_text_box(slide, Inches(7.0), Inches(2.0), Inches(5.6), Inches(0.4), "18 Edge Kinds", 17, True, MEDIUM_BLUE)

add_text_box(slide, Inches(7.0), Inches(2.5), Inches(5.6), Inches(0.3), "11 Structural (existence, ownership, containment):", 13, True, DARK_GRAY)
struct_items = ["MODULE_DECLARES_*, MODULE_IMPORTS/EXPORTS, MODULE_PROVIDES_SERVICE",
    "ROUTE_HAS_CHILD, ROUTE_ACTIVATES_COMPONENT",
    "COMPONENT_CONTAINS_WIDGET, WIDGET_CONTAINS_WIDGET",
    "COMPONENT_COMPOSES_COMPONENT"]
add_bullet_list(slide, Inches(7.0), Inches(2.8), Inches(5.6), Inches(1.5), struct_items, 11, MED_GRAY)

add_text_box(slide, Inches(7.0), Inches(4.1), Inches(5.6), Inches(0.3), "7 Executable (what can happen at runtime):", 13, True, DARK_GRAY)
exec_items = ["WIDGET_NAVIGATES_ROUTE / EXTERNAL",
    "WIDGET_TRIGGERS_HANDLER, WIDGET_SUBMITS_FORM",
    "COMPONENT_CALLS_SERVICE, COMPONENT_NAVIGATES_ROUTE",
    "ROUTE_REDIRECTS_TO_ROUTE (system)"]
add_bullet_list(slide, Inches(7.0), Inches(4.4), Inches(5.6), Inches(1.5), exec_items, 11, MED_GRAY)

# Key properties
add_text_box(slide, Inches(7.0), Inches(5.6), Inches(5.6), Inches(0.3), "Key properties:", 13, True, DARK_GRAY)
key_items = ["Every node/edge backed by SourceRef (file, char offsets)",
    "ConstraintSurface on every executable edge (guards, params, UI atoms)",
    "Deterministic: same codebase \u2192 byte-identical output"]
add_bullet_list(slide, Inches(7.0), Inches(5.9), Inches(5.6), Inches(1.0), key_items, 11, MED_GRAY)

# Stats box
add_text_box(slide, Inches(0.4), Inches(5.2), Inches(6.0), Inches(0.4),
    "Across 6 subjects: 665 nodes, 1345 edges, 993 structural, 352 executable", 14, True, MEDIUM_BLUE)

# ============================================================
# SLIDE 8: A2 — Task Workflows
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_header(slide, "A2: Task Workflow Enumeration and Classification (RQ1, RQ2)")

add_text_box(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(0.5),
    "A2 enumerates exactly one TaskWorkflow per trigger edge, producing the workflow space W.",
    17, False, DARK_GRAY)

# Left: algorithm
add_shape_bg(slide, Inches(0.4), Inches(1.9), Inches(6.0), Inches(3.5), LIGHT_GRAY)
add_text_box(slide, Inches(0.6), Inches(2.0), Inches(5.6), Inches(0.4), "Enumeration Algorithm", 17, True, MEDIUM_BLUE)
algo_items = [
    "For each component-bearing route r, compute active widgets",
    "For each enabled trigger edge (WTH/WSF/WNR/WNE):",
    "  \u2022 Collect handler-scoped effects (CCS, CNR) via effectGroupId",
    "  \u2022 Apply deterministic redirect closure (cycle-safe)",
    "  \u2022 Determine terminal node (Route or External)",
    "Aggregate: same trigger active on N routes \u2192 1 workflow",
    "Merge constraints across all steps: C(w) = \u222a constraints(e_i)",
]
add_bullet_list(slide, Inches(0.6), Inches(2.5), Inches(5.6), Inches(2.8), algo_items, 13)

# Right: classification
add_shape_bg(slide, Inches(6.8), Inches(1.9), Inches(6.0), Inches(3.5), RGBColor(0xE8, 0xF0, 0xFE))
add_text_box(slide, Inches(7.0), Inches(2.0), Inches(5.6), Inches(0.4), "Classification Rules", 17, True, MEDIUM_BLUE)
class_items = [
    "FEASIBLE: no guards, no params, no contradictions",
    "CONDITIONAL: requires runtime assignment (guards, params, form validity, expression-bound UI gates)",
    "PRUNED: provable literal contradictions only (mutex, exclusive roles, visible=false, min>max, redirect deadlock)",
    "Strict order: PRUNED checks \u2192 CONDITIONAL checks \u2192 FEASIBLE",
    "No full SAT, no heuristic tie-breaking",
]
add_bullet_list(slide, Inches(7.0), Inches(2.5), Inches(5.6), Inches(2.8), class_items, 13)

# Bottom: stats table
wf_data = [
    ["Subject", "Workflows", "FEASIBLE", "CONDITIONAL", "PRUNED", "Trigger Edges", "Routes"],
    ["posts-users-ui-ng", "18", "12", "6", "0", "18", "7"],
    ["spring-petclinic-angular", "74", "40", "34", "0", "74", "22"],
    ["heroes-angular", "19", "19", "0", "0", "19", "4"],
    ["softscanner-cqa-frontend", "16", "15", "1", "0", "16", "1"],
    ["ever-traduora", "109", "46", "63", "0", "109", "18"],
    ["airbus-inventory", "21", "13", "8", "0", "21", "6"],
    ["Total", "257", "145", "112", "0", "257", "58"],
]
add_table(slide, Inches(0.4), Inches(5.6), Inches(12.5), Inches(1.8), 8, 7, wf_data,
    [Inches(2.5), Inches(1.3), Inches(1.3), Inches(1.8), Inches(1.3), Inches(1.8), Inches(1.3)], font_size=11)

# ============================================================
# SLIDE 9: Phase B Artifact Chain
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_header(slide, "Phase B: From Workflows to Executable Tests (RQ3)")

# SubjectManifest
add_shape_bg(slide, Inches(0.4), Inches(1.3), Inches(4.0), Inches(2.8), RGBColor(0xFD, 0xF0, 0xE0))
add_text_box(slide, Inches(0.5), Inches(1.35), Inches(3.8), Inches(0.35), "B0: SubjectManifest", 15, True, ACCENT_ORANGE)
b0_items = [
    "Per-subject manual config: baseUrl, accounts, routeParamValues",
    "Maps guard names \u2192 credentials via guardSatisfies",
    "authSetup: login mechanism (route, selectors)",
    "formDataOverrides, skipWorkflows (optional)",
    "executionConfig: readiness check, seedDataNotes",
    "Validated against A2 (guards, params, workflow IDs)",
]
add_bullet_list(slide, Inches(0.5), Inches(1.75), Inches(3.8), Inches(2.2), b0_items, 11, DARK_GRAY)

# RealizationIntent
add_shape_bg(slide, Inches(4.7), Inches(1.3), Inches(4.0), Inches(2.8), RGBColor(0xDE, 0xEA, 0xF7))
add_text_box(slide, Inches(4.8), Inches(1.35), Inches(3.8), Inches(0.35), "B1.1: RealizationIntent", 15, True, MEDIUM_BLUE)
b11_items = [
    "1 intent per TaskWorkflow (257 total)",
    "Trigger widget: tagName, attributes, locator info",
    "Start route with URL template + required params",
    "Form schema (for WSF triggers): field types, constraints",
    "Guard names, param requirements, unresolved targets",
    "Deterministic derivation from A1 + A2 (no manifest)",
]
add_bullet_list(slide, Inches(4.8), Inches(1.75), Inches(3.8), Inches(2.2), b11_items, 11, DARK_GRAY)

# ActionPlan
add_shape_bg(slide, Inches(9.0), Inches(1.3), Inches(4.0), Inches(2.8), RGBColor(0xDE, 0xEA, 0xF7))
add_text_box(slide, Inches(9.1), Inches(1.35), Inches(3.8), Inches(0.35), "B1.2: ActionPlan", 15, True, MEDIUM_BLUE)
b12_items = [
    "Binds intent to manifest: account, route params, form data",
    "PreConditions: auth-setup, navigate-to-route, dialog-open",
    "ActionSteps: click, type, submit, select-option, wait-*",
    "ScopedLocator: component + form scoping, 10 strategies",
    "PostConditions: assert-url-matches, assert-no-crash",
    "257 plans, all GT-validated (0 mismatches)",
]
add_bullet_list(slide, Inches(9.1), Inches(1.75), Inches(3.8), Inches(2.2), b12_items, 11, DARK_GRAY)

# B2 Generated Tests
add_shape_bg(slide, Inches(0.4), Inches(4.4), Inches(6.0), Inches(3.0), RGBColor(0xE8, 0xF5, 0xE9))
add_text_box(slide, Inches(0.5), Inches(4.45), Inches(5.8), Inches(0.35), "B2: Generated Selenium Tests", 15, True, GREEN)
b2_items = [
    "1 .test.ts file per ActionPlan (257 total across 6 subjects)",
    "Selenium WebDriver TypeScript: Chrome headless setup, implicit wait",
    "Preconditions \u2192 Steps \u2192 Postconditions \u2192 driver.quit()",
    "Locator translation: ScopedLocator \u2192 Selenium By.* expression",
    "Component + form scoping via nested findElement chains",
    "Deterministic: same ActionPlan \u2192 byte-identical test code",
    "Generation coverage: 257/257 = 100% across all 6 subjects",
]
add_bullet_list(slide, Inches(0.5), Inches(4.85), Inches(5.8), Inches(2.5), b2_items, 12, DARK_GRAY)

# B3/B4 planned
add_shape_bg(slide, Inches(6.8), Inches(4.4), Inches(6.0), Inches(3.0), RGBColor(0xFF, 0xF8, 0xE1))
add_text_box(slide, Inches(7.0), Inches(4.45), Inches(5.6), Inches(0.35), "B3/B4: Execution + Coverage (Planned)", 15, True, AMBER)
b34_items = [
    "B3: Execute tests against running subject apps",
    "Bounded retry: 3 levels (as-is, extended wait, LLM-repair)",
    "ExecutionResult with 7-outcome taxonomy (PASS through FAIL_UNKNOWN)",
    "FAIL_APP_NOT_READY excluded from C3 denominator",
    "B4: Tiered coverage: C1 (plan), C2 (code), C3 (execution), C4 (oracle)",
    "Entity provisioning: user responsibility (seedDataNotes)",
    "NOT YET VALIDATED — execution-readiness pass complete, B3 pending",
]
add_bullet_list(slide, Inches(7.0), Inches(4.85), Inches(5.6), Inches(2.5), b34_items, 12, DARK_GRAY)


# ============================================================
# SLIDE 10: Ground Truth Construction
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_header(slide, "Ground-Truth Construction and Alignment")

add_text_box(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(0.5),
    "257 manually constructed ground-truth entries across 6 subjects align B1 outputs with expert expectations.",
    17, False, DARK_GRAY)

# Why GT matters
add_shape_bg(slide, Inches(0.4), Inches(1.9), Inches(5.8), Inches(2.5), LIGHT_GRAY)
add_text_box(slide, Inches(0.6), Inches(2.0), Inches(5.4), Inches(0.4), "Why Ground Truth Matters", 17, True, MEDIUM_BLUE)
gt_items = [
    "Validates that derived intents and plans match human understanding of the application",
    "Catches systematic errors in locator derivation, start route selection, auth materialization",
    "Identifies spec gaps that required normative amendments (5 frozen in approach.md)",
    "Prevents silent drift between implementation and spec over 257 workflows",
]
add_bullet_list(slide, Inches(0.6), Inches(2.5), Inches(5.4), Inches(1.8), gt_items, 13)

# GT process
add_shape_bg(slide, Inches(6.6), Inches(1.9), Inches(6.2), Inches(2.5), RGBColor(0xE8, 0xF0, 0xFE))
add_text_box(slide, Inches(6.8), Inches(2.0), Inches(5.8), Inches(0.4), "GT-Driven Adjudication", 17, True, MEDIUM_BLUE)
adj_items = [
    "150 initial mismatches across 6 subjects",
    "~100 GT errors (over-specification) \u2192 GT repaired per normative rules",
    "5 spec gaps \u2192 normative amendments to approach.md",
    "3 validator issues \u2192 comparison logic fixed",
    "2 implementation bugs \u2192 code fixed",
    "Final: 257/257 matched for both B1 intents and B1 plans",
]
add_bullet_list(slide, Inches(6.8), Inches(2.5), Inches(5.8), Inches(1.8), adj_items, 13)

# GT repairs table
repair_data = [
    ["Subject", "GT Entries", "Repairs", "Final Match"],
    ["posts-users-ui-ng", "18", "7", "18/18"],
    ["heroes-angular", "19", "4", "19/19"],
    ["softscanner-cqa-frontend", "16", "2", "16/16"],
    ["spring-petclinic-angular", "74", "116", "74/74"],
    ["ever-traduora", "109", "123", "109/109"],
    ["airbus-inventory", "21", "16", "21/21"],
    ["Total", "257", "266", "257/257"],
]
add_table(slide, Inches(0.4), Inches(4.7), Inches(8.0), Inches(2.5), 8, 4, repair_data,
    [Inches(2.8), Inches(1.5), Inches(1.2), Inches(1.5)], font_size=12)

# ============================================================
# SLIDE 11: Validation Corpus
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_header(slide, "Validation Corpus: 6 Angular Applications")

corpus_data = [
    ["Subject", "Framework", "Nodes", "Edges", "Workflows", "Auth", "Key Characteristics"],
    ["posts-users-ui-ng", "Angular 18", "64", "131", "18", "No", "Social platform, forms, route params"],
    ["spring-petclinic-angular", "Angular 14", "191", "422", "74", "No", "Healthcare CRUD, backend deps, params"],
    ["heroes-angular", "Angular 14", "67", "88", "19", "No", "NgRx, external links, lazy modules"],
    ["softscanner-cqa-frontend", "Angular 17", "37", "74", "16", "No", "Single route, dialogs, composition"],
    ["ever-traduora", "Angular 12", "247", "499", "109", "Yes", "Auth guards, 20 routes, largest graph"],
    ["airbus-inventory", "Angular 12", "59", "131", "21", "Yes", "JWT guards, flat module, Spring Boot"],
    ["Total", "", "665", "1345", "257", "", "Angular 12\u201318, varied architecture"],
]
add_table(slide, Inches(0.4), Inches(1.3), Inches(12.5), Inches(3.2), 8, 7, corpus_data,
    [Inches(2.5), Inches(1.3), Inches(0.8), Inches(0.8), Inches(1.3), Inches(0.6), Inches(5.0)], font_size=12)

add_text_box(slide, Inches(0.6), Inches(4.8), Inches(12), Inches(0.5),
    "Selection criteria: varied architecture (flat vs. modular), Angular versions 12\u201318, with/without auth guards, different backend stacks.", 14, False, MED_GRAY)

# Diversity dimensions
div_items = [
    "Module structure: flat (airbus: 1 module) to highly modular (traduora: 4 modules, 20 routes, 45 components)",
    "Auth patterns: no auth (4 subjects), route guards with roles (2 subjects: traduora, airbus)",
    "UI patterns: standard forms, Angular Material, NgRx, dialog/modal composition, external navigation",
    "Backend dependencies: REST APIs, JSON Server, in-memory, Spring Boot + MySQL, Express.js",
    "All subjects validated for byte-identical determinism across repeated runs",
]
add_bullet_list(slide, Inches(0.6), Inches(5.3), Inches(12), Inches(2.0), div_items, 13)

# ============================================================
# SLIDE 12: Current Validated Results
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_header(slide, "Current Validated Results")

# Results table
results_data = [
    ["Stage", "Metric", "Result", "Status"],
    ["A1", "Multigraph extraction (6 subjects)", "665 nodes, 1345 edges", "\u2713 Complete"],
    ["A2", "Workflow enumeration", "257 workflows (145F + 112C + 0P)", "\u2713 Complete"],
    ["B0", "Manifest validation", "6/6 VALID, 0 errors", "\u2713 Complete"],
    ["B1.1", "RealizationIntent derivation", "257/257 GT match", "\u2713 Complete"],
    ["B1.2", "ActionPlan generation", "257/257 GT match", "\u2713 Complete"],
    ["B2", "Selenium test generation", "257/257 generated (100%)", "\u2713 Complete"],
    ["B3", "Test execution at scale", "Pending", "\u25cb Planned"],
    ["B4", "Coverage reporting (C1\u2013C4)", "Pending", "\u25cb Planned"],
]
add_table(slide, Inches(0.4), Inches(1.3), Inches(12.5), Inches(3.6), 9, 4, results_data,
    [Inches(1.2), Inches(4.5), Inches(4.5), Inches(2.0)], font_size=13)

# Determinism box
add_shape_bg(slide, Inches(0.4), Inches(5.2), Inches(6.0), Inches(2.0), RGBColor(0xE8, 0xF5, 0xE9))
add_text_box(slide, Inches(0.6), Inches(5.3), Inches(5.6), Inches(0.4), "Determinism Verified", 16, True, GREEN)
det_items = [
    "\u2713 A1 + A2 bundles byte-identical across runs",
    "\u2713 B0 summary byte-identical",
    "\u2713 B1 intents + plans byte-identical",
    "\u2713 B2 summary byte-identical",
]
add_bullet_list(slide, Inches(0.6), Inches(5.7), Inches(5.6), Inches(1.3), det_items, 13, GREEN)

# Test suite box
add_shape_bg(slide, Inches(6.8), Inches(5.2), Inches(6.0), Inches(2.0), LIGHT_GRAY)
add_text_box(slide, Inches(7.0), Inches(5.3), Inches(5.6), Inches(0.4), "Test Suite", 16, True, MEDIUM_BLUE)
test_items = [
    "248 automated tests (unit + integration)",
    "Typecheck (source + tests), lint — all green",
    "All 6 subjects validated for every stage",
    "GT validation: 257/257 for intents and plans",
]
add_bullet_list(slide, Inches(7.0), Inches(5.7), Inches(5.6), Inches(1.3), test_items, 13)

# ============================================================
# SLIDE 13: Coverage Model
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_header(slide, "Coverage Model: Generation vs. Execution vs. Oracle")

# Coverage tiers
tier_data = [
    ["Tier", "Definition", "Denominator", "Current Status"],
    ["C1 (Plan)", "Fraction of W with a valid ActionPlan", "|W| (non-PRUNED)", "257/257 = 100%"],
    ["C2 (Code)", "Fraction of W with syntactically valid test code", "|W|", "257/257 = 100%"],
    ["C3 (Execution)", "Fraction of W with a passing test", "|W| \u2212 FAIL_APP_NOT_READY", "Pending (B3)"],
    ["C4 (Oracle)", "Assertion richness beyond URL-match", "Deferred", "Deferred"],
]
add_table(slide, Inches(0.4), Inches(1.3), Inches(12.5), Inches(2.2), 5, 4, tier_data,
    [Inches(1.8), Inches(4.5), Inches(3.0), Inches(3.0)], font_size=13)

# Key distinctions
add_shape_bg(slide, Inches(0.4), Inches(3.8), Inches(6.0), Inches(3.4), LIGHT_GRAY)
add_text_box(slide, Inches(0.6), Inches(3.9), Inches(5.6), Inches(0.4),
    "Generation Coverage (C1+C2) \u2014 Validated", 16, True, GREEN)
gen_items = [
    "Measures: can we plan and generate code for every workflow?",
    "Current: 257/257 = 100% for both C1 and C2",
    "This is a necessary but not sufficient condition for E2E success",
    "Does NOT prove that generated tests will execute correctly",
]
add_bullet_list(slide, Inches(0.6), Inches(4.4), Inches(5.6), Inches(2.5), gen_items, 13)

add_shape_bg(slide, Inches(6.8), Inches(3.8), Inches(6.0), Inches(3.4), RGBColor(0xFF, 0xF8, 0xE1))
add_text_box(slide, Inches(7.0), Inches(3.9), Inches(5.6), Inches(0.4),
    "Execution Coverage (C3) \u2014 Pending", 16, True, AMBER)
exec_items = [
    "Measures: do generated tests actually pass on a running app?",
    "Requires: running application, seeded data, correct auth",
    "125 workflows need no entity data (49%)",
    "129 require entity data in database (50%)",
    "97 require authentication (38%)",
    "Execution-readiness audit complete; B3 implementation pending",
]
add_bullet_list(slide, Inches(7.0), Inches(4.4), Inches(5.6), Inches(2.5), exec_items, 13)


# ============================================================
# SLIDE 14: Comparison with AutoE2E
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_header(slide, "Comparison with AutoE2E")

add_text_box(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(0.5),
    "AutoE2E (Allan et al., ICSE 2025): LLM-driven feature discovery via runtime exploration.",
    16, False, DARK_GRAY)

# Comparison table
comp_data = [
    ["Dimension", "Our Approach", "AutoE2E"],
    ["Interaction space", "Closed workflow space W from source code", "Open: discovered at runtime via crawler"],
    ["Coverage denominator", "Fixed |W| (implementation-grounded)", "Exploration-bounded (30-action ceiling)"],
    ["Constraint handling", "Explicit: guards, params, UI atoms in C(w)", "Implicit: auth wall blocks navigation"],
    ["Test generation", "Deterministic Selenium from ActionPlans", "Does not generate executable tests *"],
    ["Coverage definition", "|E|/|W| over same fixed W for both", "Feature discovery count (no W defined)"],
    ["Auth support", "Manifest-driven credentials + authSetup", "No credential management"],
    ["Form handling", "Schema-aware from A1 extraction", "Requires data-testid (PetClinic only)"],
    ["Determinism", "Byte-identical across runs", "LLM-dependent, non-deterministic"],
]
add_table(slide, Inches(0.4), Inches(1.8), Inches(12.5), Inches(3.4), 9, 3, comp_data,
    [Inches(2.2), Inches(5.2), Inches(5.2)], font_size=12)

add_text_box(slide, Inches(0.6), Inches(5.3), Inches(12), Inches(0.3),
    "* AutoE2E produces a state graph and feature database, not executable test files. Confirmed by code inspection and 6 benchmark runs.",
    12, False, MED_GRAY)

# Paper results
add_shape_bg(slide, Inches(0.4), Inches(5.7), Inches(12.5), Inches(1.5), RGBColor(0xFD, 0xF0, 0xE0))
add_text_box(slide, Inches(0.6), Inches(5.75), Inches(12), Inches(0.4),
    "Paper Evaluation (2 subjects, same W, execution-validated):", 15, True, ACCENT_ORANGE)
paper_items = [
    "PetClinic: |W|=84, Ours 71/84 (84.5%) vs AutoE2E 49/84 (58.3%) \u2014 mapped to same action alphabet",
    "Tour of Heroes: |W|=36, Ours 32/36 (88.9%) vs AutoE2E 23/36 (63.9%) \u2014 mapped to same action alphabet",
    "Caveat: these are execution-validated results from the paper's evaluation; current 6-subject execution coverage is pending B3",
]
add_bullet_list(slide, Inches(0.6), Inches(6.15), Inches(12), Inches(1.0), paper_items, 12, DARK_GRAY, Pt(3))


# ============================================================
# SLIDE 15: AutoE2E Benchmark Evidence
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_header(slide, "AutoE2E: 6-Subject Benchmark Evidence")

add_text_box(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(0.5),
    "Independent evaluation: AutoE2E pipeline executed on all 6 subjects used in our validation corpus.",
    16, False, DARK_GRAY)

ae_data = [
    ["Subject", "States", "Actions\n(of 30)", "Features\nStored", "Forms\nSubmitted", "Binding\nConstraint"],
    ["PetClinic", "9", "30", "46", "4*", "30-action ceiling"],
    ["Ever Traduora", "3", "13", "5", "0", "Auth wall"],
    ["Posts & Users", "7", "30", "~37", "0", "30-action ceiling"],
    ["Heroes Angular", "5", "30", "~25", "0", "30-action ceiling"],
    ["Airbus Inventory", "1", "6", "0", "0", "Auth wall + finality"],
    ["SoftScanner CQA", "9", "30", "6", "0", "30-action ceiling"],
]
add_table(slide, Inches(0.4), Inches(1.8), Inches(12.5), Inches(2.8), 7, 6, ae_data,
    [Inches(2.5), Inches(1.0), Inches(1.5), Inches(1.5), Inches(1.5), Inches(4.5)], font_size=12)

add_text_box(slide, Inches(0.6), Inches(4.7), Inches(12), Inches(0.3),
    "* PetClinic forms use custom data-testid attributes; mechanical traversal only, business outcomes not verified.",
    12, False, MED_GRAY)

# Key findings
add_shape_bg(slide, Inches(0.4), Inches(5.1), Inches(12.5), Inches(2.2), LIGHT_GRAY)
add_text_box(slide, Inches(0.6), Inches(5.15), Inches(12), Inches(0.4),
    "Key Findings from AutoE2E Evaluation", 16, True, MEDIUM_BLUE)
ae_findings = [
    "AutoE2E does not generate executable test files — produces state graph + feature database only",
    "Form submission requires PetClinic-specific data-testid attributes; fails on all standard Angular forms",
    "Auth-protected apps (traduora, airbus) severely constrained: pipeline has no credential management",
    "30-action ceiling is primary bottleneck for accessible apps; does not scale to large interaction surfaces",
    "Comparison is fundamentally asymmetric: our W is grounded in source code; AutoE2E's space is exploration-bounded",
]
add_bullet_list(slide, Inches(0.6), Inches(5.6), Inches(12), Inches(1.6), ae_findings, 13, DARK_GRAY, Pt(3))


# ============================================================
# SLIDE 16: Limitations and Open Issues
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_header(slide, "Current Limitations and Open Issues")

# Left: execution readiness
add_shape_bg(slide, Inches(0.4), Inches(1.3), Inches(6.0), Inches(5.8), RGBColor(0xFF, 0xF8, 0xE1))
add_text_box(slide, Inches(0.6), Inches(1.4), Inches(5.6), Inches(0.4),
    "Execution Readiness (Not Yet Validated)", 17, True, AMBER)
exec_lim = [
    "B3 not implemented: generated tests not yet executed at scale",
    "129 workflows (50%) depend on entity data existing in the running app",
    "Entity provisioning is the user's responsibility (seedDataNotes)",
    "No automated readiness verification beyond HTTP GET check",
    "Failure taxonomy defined but not empirically validated",
    "LLM-assisted retry (Level 3) not yet implemented",
    "C3 (execution coverage) remains the key open metric",
]
add_bullet_list(slide, Inches(0.6), Inches(1.9), Inches(5.6), Inches(4.8), exec_lim, 14)

# Right: methodological
add_shape_bg(slide, Inches(6.8), Inches(1.3), Inches(6.0), Inches(5.8), LIGHT_GRAY)
add_text_box(slide, Inches(7.0), Inches(1.4), Inches(5.6), Inches(0.4),
    "Methodological Limitations", 17, True, MEDIUM_BLUE)
meth_lim = [
    "Angular-specific: AST extraction uses ts-morph + @angular/compiler",
    "Bounded path length: W depends on exploration depth k=5",
    "Single-trigger model: no multi-step workflow composition",
    "Postconditions limited to URL-match + no-crash (C4 deferred)",
    "No @Output() tracing, no ancestor *ngIf propagation",
    "Paper evaluation on 2 subjects only; 6-subject execution pending",
    "Mapping AutoE2E actions to our W is deterministic but approximate",
]
add_bullet_list(slide, Inches(7.0), Inches(1.9), Inches(5.6), Inches(4.8), meth_lim, 14)


# ============================================================
# SLIDE 17: Conclusion and Contributions
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
section_header(slide, "Conclusion and Contributions")

add_text_box(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(0.5),
    "This work reframes E2E test coverage as a property of a finite, constraint-aware workflow space grounded in the frontend implementation.",
    17, True, DARK_BLUE)

# Contributions
add_shape_bg(slide, Inches(0.4), Inches(2.0), Inches(6.0), Inches(4.0), LIGHT_GRAY)
add_text_box(slide, Inches(0.6), Inches(2.1), Inches(5.6), Inches(0.4), "Contributions", 18, True, MEDIUM_BLUE)
contrib_items = [
    "Augmented workflows: navigation paths annotated with implementation-level guard/state constraints",
    "Finite, closed workflow space W as the coverage denominator — implementation-grounded, not exploration-bounded",
    "Automated pipeline: A1 (multigraph) \u2192 A2 (workflows) \u2192 B0\u2013B2 (realization + code generation)",
    "Deterministic throughout: same codebase \u2192 byte-identical artifacts",
    "GT-validated on 6 Angular subjects (257 workflows, 0 mismatches)",
    "Separation of construction (Phase A) from realization (Phase B) — coverage measures realization effectiveness",
]
add_bullet_list(slide, Inches(0.6), Inches(2.5), Inches(5.6), Inches(3.3), contrib_items, 14)

# What it shows
add_shape_bg(slide, Inches(6.8), Inches(2.0), Inches(6.0), Inches(4.0), RGBColor(0xE8, 0xF5, 0xE9))
add_text_box(slide, Inches(7.0), Inches(2.1), Inches(5.6), Inches(0.4), "What the Evidence Shows", 18, True, GREEN)
evidence_items = [
    "C1 + C2 = 100%: every workflow in W has a valid plan and generated test",
    "Paper evaluation (2 subjects): 84.5% and 88.9% execution coverage vs 58.3% and 63.9% for AutoE2E",
    "Consistent improvement on applications with guarded routes and parameterized workflows",
    "Coverage differences reflect realization capability, not denominator differences (same W)",
]
add_bullet_list(slide, Inches(7.0), Inches(2.5), Inches(5.6), Inches(2.0), evidence_items, 14, DARK_GRAY)

# Future work
add_shape_bg(slide, Inches(6.8), Inches(4.3), Inches(6.0), Inches(1.7), RGBColor(0xFF, 0xF8, 0xE1))
add_text_box(slide, Inches(7.0), Inches(4.35), Inches(5.6), Inches(0.4), "Future Work", 16, True, AMBER)
future_items = [
    "B3: Execute generated tests at scale across all 6 subjects",
    "Extend to additional frontend frameworks (React, Vue)",
    "Richer postconditions and oracle mechanisms (C4)",
]
add_bullet_list(slide, Inches(7.0), Inches(4.7), Inches(5.6), Inches(1.2), future_items, 13)

# Bottom: honest caveat
add_shape_bg(slide, Inches(0.4), Inches(6.3), Inches(12.5), Inches(0.7), RGBColor(0xFD, 0xF0, 0xE0))
add_text_box(slide, Inches(0.6), Inches(6.35), Inches(12), Inches(0.6),
    "Honest status: Generation coverage validated (100%). Execution coverage validated on 2 subjects in paper. Full 6-subject execution coverage pending B3 implementation.",
    14, True, ACCENT_ORANGE, PP_ALIGN.CENTER)

# ============================================================
# Save
# ============================================================
output_path = os.path.join(os.path.dirname(__file__), "phase-b-overview.pptx")
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
